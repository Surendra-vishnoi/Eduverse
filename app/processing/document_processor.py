import asyncio
import os
import tempfile
from pathlib import Path
from typing import List, Optional

from langchain_core.documents import Document

SUPPORTED_DOCUMENT_FORMATS = {".docx", ".pptx", ".txt", ".md"}


def _extract_docx_text(file_path: str) -> str:
    from docx import Document as DocxDocument

    doc = DocxDocument(file_path)
    chunks: list[str] = []

    for paragraph in doc.paragraphs:
        text = (paragraph.text or "").strip()
        if text:
            chunks.append(text)

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join((cell.text or "").strip() for cell in row.cells)
            if row_text.strip(" |"):
                chunks.append(row_text)

    return "\n".join(chunks)


def _extract_pptx_slides(file_path: str) -> list[tuple[int, str]]:
    from pptx import Presentation

    presentation = Presentation(file_path)
    slides: list[tuple[int, str]] = []

    for idx, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            text = (text or "").strip()
            if text:
                slide_lines.append(text)

        combined = "\n".join(slide_lines).strip()
        if combined:
            slides.append((idx, combined))

    return slides


def _extract_text_like(file_path: str) -> str:
    return Path(file_path).read_text(encoding="utf-8", errors="ignore")


def _process_document_sync(
    file_path: str,
    file_name: str,
    course_id: Optional[str],
    source_id: Optional[str],
) -> List[Document]:
    ext = Path(file_name).suffix.lower()
    source = source_id or file_name

    if ext == ".docx":
        text = _extract_docx_text(file_path).strip()
        if not text:
            return []
        return [
            Document(
                page_content=text,
                metadata={
                    "source_type": "document",
                    "source_id": source,
                    "file_name": file_name,
                    "course_id": course_id,
                    "page_number": 1,
                    "start_time": None,
                    "end_time": None,
                    "contains_visual": False,
                },
            )
        ]

    if ext == ".pptx":
        docs: list[Document] = []
        for slide_no, text in _extract_pptx_slides(file_path):
            cleaned = text.strip()
            if not cleaned:
                continue
            docs.append(
                Document(
                    page_content=f"[SLIDE {slide_no}]\n{cleaned}",
                    metadata={
                        "source_type": "document",
                        "source_id": source,
                        "file_name": file_name,
                        "course_id": course_id,
                        "page_number": slide_no,
                        "start_time": None,
                        "end_time": None,
                        "contains_visual": False,
                    },
                )
            )
        return docs

    if ext in {".txt", ".md"}:
        text = _extract_text_like(file_path).strip()
        if not text:
            return []
        return [
            Document(
                page_content=text,
                metadata={
                    "source_type": "document",
                    "source_id": source,
                    "file_name": file_name,
                    "course_id": course_id,
                    "page_number": 1,
                    "start_time": None,
                    "end_time": None,
                    "contains_visual": False,
                },
            )
        ]

    return []


async def process_document(
    file_path: str,
    file_name: str,
    course_id: Optional[str] = None,
    source_id: Optional[str] = None,
) -> List[Document]:
    return await asyncio.to_thread(
        _process_document_sync,
        file_path,
        file_name,
        course_id,
        source_id,
    )


async def process_document_bytes(
    file_bytes: bytes,
    file_name: str,
    course_id: Optional[str] = None,
    source_id: Optional[str] = None,
) -> List[Document]:
    ext = os.path.splitext(file_name)[1] or ".txt"
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        return await process_document(tmp_path, file_name, course_id, source_id)
    finally:
        os.unlink(tmp_path)